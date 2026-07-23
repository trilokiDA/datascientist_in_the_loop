"""
Create a comprehensive PowerPoint presentation for the EDA Pipeline project
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from datetime import datetime

def create_title_slide(prs):
    """Create title slide"""
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "EDA Pipeline with Agentic Workflow"
    subtitle.text = "Intelligent, Production-Ready Exploratory Data Analysis\nPowered by LangGraph, Groq, and Open-Source Tools\n\nVersion 3.3 - July 2026"

    # Style title
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 119, 180)

    # Style subtitle
    for paragraph in subtitle.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.alignment = PP_ALIGN.CENTER

def create_agenda_slide(prs):
    """Create agenda slide"""
    slide_layout = prs.slide_layouts[1]  # Title and content
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "Agenda"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    agenda_items = [
        "Project Overview",
        "Key Features & Latest Updates",
        "Architecture & Tech Stack",
        "6 Specialized Agents",
        "Human-in-the-Loop Approval Gates",
        "Transformation & Export Capabilities",
        "Demo & Use Cases",
        "Future Roadmap"
    ]

    for item in agenda_items:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(22)

def create_overview_slide(prs):
    """Create project overview slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "Project Overview"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = "What is EDA Pipeline?"
    p.font.size = Pt(24)
    p.font.bold = True
    p.space_after = Pt(12)

    overview_text = [
        "An intelligent, production-ready EDA pipeline powered by AI agents",
        "Automated exploratory data analysis with explainable AI",
        "Complete suite of 6 specialized agents for comprehensive analysis",
        "Interactive Streamlit UI with real-time progress tracking",
        "Supports CSV and Excel files (.xlsx, .xls)",
        "Hybrid scaling: handles small to large datasets intelligently",
        "Human-in-the-loop workflow with approval gates",
        "Export capabilities: HTML reports, JSON data, transformed CSV"
    ]

    for text in overview_text:
        p = tf.add_paragraph()
        p.text = text
        p.level = 0
        p.font.size = Pt(18)
        p.space_after = Pt(8)

def create_key_features_slide(prs):
    """Create key features slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "Key Features - Version 3.3"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    features = [
        ("NEW in v3.3: Excel File Support", [
            "Support for .xlsx and .xls files",
            "Automatic detection and loading of first sheet",
            "Seamless integration with all agents"
        ]),
        ("Human-in-the-Loop Approval Gates (v3.2)", [
            "Review and approve each agent before proceeding",
            "4 decision options: Approve, Retry, Skip, Stop",
            "Complete decision history tracking"
        ]),
        ("Multi-Transformation Selection (v3.1)", [
            "Select and apply multiple transformations at once",
            "Preview combined effects before applying",
            "Export transformed CSV for downstream use"
        ])
    ]

    for feature_name, details in features:
        p = tf.add_paragraph()
        p.text = feature_name
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(31, 119, 180)
        p.space_after = Pt(8)

        for detail in details:
            p = tf.add_paragraph()
            p.text = detail
            p.level = 1
            p.font.size = Pt(16)
            p.space_after = Pt(4)

        p.space_after = Pt(12)

def create_architecture_slide(prs):
    """Create architecture slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "System Architecture"

    # Add text box for architecture description
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(5)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame

    arch_text = """
┌─────────────────┐
│  Streamlit UI   │  ← Interactive interface with progress tracking
└────────┬────────┘
         │
┌────────▼────────┐
│  LangGraph      │  ← State machine orchestration
│  Orchestration  │
└────────┬────────┘
         │
┌────────▼────────────────────────────────────────────┐
│              Specialized Agents                      │
├──────────┬──────────┬──────────┬──────────┬─────────┤
│ Profile  │ Quality  │Transform │Visualize │Features │
│  Agent   │  Agent   │  Agent   │  Agent   │ Agent   │
│          │          │          │          │  Stat   │
└──────────┴──────────┴──────────┴──────────┴─────────┘
         │
┌────────▼────────┐
│   Data Layer    │  ← Pandas/DuckDB with intelligent switching
└─────────────────┘
    """

    p = tf.add_paragraph()
    p.text = arch_text
    p.font.name = 'Courier New'
    p.font.size = Pt(14)

def create_agents_overview_slide(prs):
    """Create agents overview slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "6 Specialized Agents"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    agents = [
        ("ProfileAgent", "Dataset profiling, structure, column types, basic statistics"),
        ("QualityAgent", "Data quality assessment, missing values, duplicates, outliers"),
        ("TransformAgent", "Data cleaning, imputation, encoding, scaling, transformations"),
        ("VisualizationAgent", "Automatic chart generation, distributions, correlations"),
        ("FeatureAgent", "Feature importance, correlation analysis, engineering suggestions"),
        ("StatAgent", "Statistical testing, normality tests, hypothesis testing")
    ]

    for agent_name, description in agents:
        p = tf.add_paragraph()
        p.text = agent_name
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(31, 119, 180)
        p.space_after = Pt(4)

        p = tf.add_paragraph()
        p.text = description
        p.level = 1
        p.font.size = Pt(16)
        p.space_after = Pt(12)

def create_profile_agent_slide(prs):
    """Detailed ProfileAgent slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "ProfileAgent - Dataset Profiling"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = "Capabilities:"
    p.font.size = Pt(22)
    p.font.bold = True

    capabilities = [
        "Dataset shape and structure analysis",
        "Column type detection (numeric, categorical, datetime)",
        "Memory usage analysis",
        "Basic statistics per column (mean, median, std, min, max)",
        "Cardinality analysis",
        "Missing value detection",
        "Issue identification (high missing, high/low cardinality)"
    ]

    for cap in capabilities:
        p = tf.add_paragraph()
        p.text = cap
        p.level = 0
        p.font.size = Pt(18)

def create_quality_agent_slide(prs):
    """Detailed QualityAgent slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "QualityAgent - Data Quality Assessment"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = "Features:"
    p.font.size = Pt(22)
    p.font.bold = True

    features = [
        "Missing value detection and pattern analysis",
        "Duplicate row identification",
        "Outlier detection using IQR method",
        "Data quality scoring",
        "Interactive visualizations:",
        "  • Missing value heatmaps",
        "  • Outlier detection box plots",
        "  • Duplicate analysis gauge charts",
        "Data type inconsistency detection"
    ]

    for feature in features:
        p = tf.add_paragraph()
        p.text = feature
        p.level = 0
        p.font.size = Pt(18)

def create_transform_agent_slide(prs):
    """Detailed TransformAgent slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "TransformAgent - Data Transformation"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = "Transformation Types:"
    p.font.size = Pt(22)
    p.font.bold = True

    transformations = [
        "Missing value handling (drop, impute median/mode/constant)",
        "Outlier handling (capping, removal)",
        "Categorical encoding (one-hot, label encoding)",
        "Numeric scaling (standard, min-max)",
        "Data type conversions",
        "Deduplication",
        "Cardinality reduction",
        "",
        "NEW Features:",
        "✓ Multi-transformation selection",
        "✓ Before/after comparison preview",
        "✓ Apply to full dataset with progress tracking",
        "✓ Export transformed CSV"
    ]

    for transform in transformations:
        p = tf.add_paragraph()
        p.text = transform
        p.level = 0
        p.font.size = Pt(16)

def create_approval_gates_slide(prs):
    """Create approval gates slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "Human-in-the-Loop Approval Gates"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = "What are Approval Gates?"
    p.font.size = Pt(22)
    p.font.bold = True

    description = [
        "Pause workflow after each agent for human review",
        "Review confidence scores, issues found, and recommendations",
        "Make informed decisions before proceeding",
        "",
        "4 Decision Options:",
        "  ✅ Approve - Continue to next agent",
        "  🔄 Retry - Re-run current agent",
        "  ⏩ Skip - Move to next agent without applying changes",
        "  ⏹️ Stop - End workflow immediately",
        "",
        "Benefits:",
        "• Critical analysis control",
        "• Compliance and auditability",
        "• Learning and understanding workflow",
        "• Complete decision history tracking"
    ]

    for line in description:
        p = tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(16)

def create_workflows_slide(prs):
    """Create workflows slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "Available Workflows"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    workflows = [
        ("Quick Analysis", "All 6 agents, comprehensive overview, ~5-10 minutes"),
        ("Deep Dive", "5 agents (no transform), thorough quality assessment"),
        ("ML Preparation", "4 agents focused on ML readiness (Profile, Quality, Feature, Transform)"),
        ("Individual Agent", "Run any single agent on-demand for specific questions")
    ]

    for workflow_name, description in workflows:
        p = tf.add_paragraph()
        p.text = workflow_name
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = RGBColor(31, 119, 180)

        p = tf.add_paragraph()
        p.text = description
        p.level = 1
        p.font.size = Pt(16)
        p.space_after = Pt(12)

    p = tf.add_paragraph()
    p.text = "All workflows available WITH or WITHOUT approval gates"
    p.font.size = Pt(18)
    p.font.italic = True

def create_tech_stack_slide(prs):
    """Create tech stack slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "Technology Stack"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    tech_stack = [
        ("Orchestration", "LangGraph - State machine workflow management"),
        ("LLM", "Groq (Llama 3.1) - Fast LLM inference"),
        ("UI", "Streamlit - Interactive web interface"),
        ("Data Processing", "Pandas, DuckDB - In-memory and large dataset handling"),
        ("Visualization", "Plotly, Seaborn, Matplotlib - Interactive charts"),
        ("Statistics", "SciPy, Statsmodels, Scikit-learn - ML prep"),
        ("Persistence", "SQLite - State checkpointing"),
        ("File Support", "openpyxl - Excel file reading")
    ]

    for category, tech in tech_stack:
        p = tf.add_paragraph()
        p.text = f"{category}: {tech}"
        p.level = 0
        p.font.size = Pt(18)

def create_export_slide(prs):
    """Create export capabilities slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "Export Capabilities"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = "Multiple Export Formats:"
    p.font.size = Pt(22)
    p.font.bold = True

    exports = [
        "📄 HTML Reports - Beautiful, interactive reports with embedded visualizations",
        "📊 JSON Data - Complete analysis results in structured format",
        "📋 Transformed CSV - Cleaned and transformed dataset for downstream use",
        "",
        "Features:",
        "• Custom export naming",
        "• Timestamp-based organization",
        "• Confidence scores and reasoning included",
        "• Shareable with stakeholders",
        "• API-ready JSON format",
        "• Full dataset processing (not just samples)"
    ]

    for export_item in exports:
        p = tf.add_paragraph()
        p.text = export_item
        p.level = 0
        p.font.size = Pt(18)

def create_use_case_slide(prs):
    """Create use case example slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "Example Use Case: Titanic Dataset"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    steps = [
        "1. Upload titanic_train.csv (891 rows, 12 columns)",
        "2. Run 'Quick Analysis' workflow",
        "3. Review results in tabbed interface",
        "4. Navigate to 'Transform' tab",
        "5. Click 'Select All High Priority' → 3 transformations selected",
        "6. Preview combined transformations:",
        "   • Sex becomes Sex_male & Sex_female (one-hot encoding)",
        "   • Missing Age values filled with median (28)",
        "   • Cabin column removed (77% missing)",
        "7. Apply 3 transformations to Full Dataset",
        "8. Export transformed CSV",
        "9. Result: 891 rows, 13 columns - ML ready! 🚀"
    ]

    for step in steps:
        p = tf.add_paragraph()
        p.text = step
        p.level = 0
        p.font.size = Pt(16)
        p.space_after = Pt(4)

def create_best_practices_slide(prs):
    """Create best practices slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "Best Practices"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    practices = [
        ("For Small Datasets (< 10k rows)", [
            "Use 'Quick Analysis' for complete overview",
            "All processing happens in-memory",
            "Fast execution (~2-5 minutes)"
        ]),
        ("For Large Datasets (> 100k rows)", [
            "System automatically switches to DuckDB backend",
            "Intelligent sampling for visualization",
            "Enable progress tracking"
        ]),
        ("For Production Use", [
            "Run Deep Dive Workflow",
            "Export HTML reports for documentation",
            "Review confidence scores before acting",
            "Use transformed CSV for downstream processing"
        ]),
        ("For ML Projects", [
            "Use ML Preparation workflow",
            "Review FeatureAgent recommendations",
            "Check StatAgent for distribution assumptions",
            "Apply and export transformations"
        ])
    ]

    for practice_name, details in practices:
        p = tf.add_paragraph()
        p.text = practice_name
        p.font.size = Pt(18)
        p.font.bold = True

        for detail in details:
            p = tf.add_paragraph()
            p.text = detail
            p.level = 1
            p.font.size = Pt(14)

def create_roadmap_slide(prs):
    """Create future roadmap slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "Future Roadmap"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = "Planned Enhancements:"
    p.font.size = Pt(22)
    p.font.bold = True

    roadmap = [
        "Additional agent types (time series, text analysis)",
        "More export formats (PDF, Excel workbooks)",
        "Advanced visualizations (3D plots, interactive dashboards)",
        "Performance optimizations for very large datasets",
        "Custom workflow builder UI",
        "Cloud deployment options",
        "API endpoints for programmatic access",
        "Integration with popular ML platforms",
        "Collaborative features (team workspaces)",
        "Scheduled/automated analysis runs"
    ]

    for item in roadmap:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(18)

def create_conclusion_slide(prs):
    """Create conclusion slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "Conclusion"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = "Key Takeaways:"
    p.font.size = Pt(24)
    p.font.bold = True
    p.space_after = Pt(12)

    takeaways = [
        "✓ Production-ready EDA pipeline with AI-powered agents",
        "✓ Comprehensive analysis with 6 specialized agents",
        "✓ Human-in-the-loop control with approval gates",
        "✓ Flexible workflows for different use cases",
        "✓ Excel and CSV support with intelligent scaling",
        "✓ Complete export capabilities for sharing and downstream use",
        "✓ Explainable AI with confidence scores and reasoning",
        "",
        "🚀 Ready to accelerate your data analysis workflow!"
    ]

    for takeaway in takeaways:
        p = tf.add_paragraph()
        p.text = takeaway
        p.level = 0
        p.font.size = Pt(18)
        p.space_after = Pt(8)

def create_thank_you_slide(prs):
    """Create thank you slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(1.5)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame

    p = tf.add_paragraph()
    p.text = "Thank You!"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(31, 119, 180)
    p.alignment = PP_ALIGN.CENTER

    # Add subtitle
    left = Inches(1)
    top = Inches(4)
    width = Inches(8)
    height = Inches(1.5)

    textbox2 = slide.shapes.add_textbox(left, top, width, height)
    tf2 = textbox2.text_frame

    p2 = tf2.add_paragraph()
    p2.text = "Questions?\n\nEDA Pipeline - Version 3.3\nJuly 2026"
    p2.font.size = Pt(24)
    p2.alignment = PP_ALIGN.CENTER

def create_presentation():
    """Main function to create the presentation"""
    print("Creating PowerPoint presentation...")

    # Create presentation object
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Create all slides
    print("Creating slides...")
    create_title_slide(prs)
    create_agenda_slide(prs)
    create_overview_slide(prs)
    create_key_features_slide(prs)
    create_architecture_slide(prs)
    create_agents_overview_slide(prs)
    create_profile_agent_slide(prs)
    create_quality_agent_slide(prs)
    create_transform_agent_slide(prs)
    create_approval_gates_slide(prs)
    create_workflows_slide(prs)
    create_tech_stack_slide(prs)
    create_export_slide(prs)
    create_use_case_slide(prs)
    create_best_practices_slide(prs)
    create_roadmap_slide(prs)
    create_conclusion_slide(prs)
    create_thank_you_slide(prs)

    # Save presentation
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"EDA_Pipeline_Presentation_{timestamp}.pptx"
    prs.save(filename)

    print(f"\nPresentation created successfully!")
    print(f"File: {filename}")
    print(f"Total slides: {len(prs.slides)}")

    return filename

if __name__ == "__main__":
    create_presentation()
